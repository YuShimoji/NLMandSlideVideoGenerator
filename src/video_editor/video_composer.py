"""
動画合成モジュール
音声、スライド、字幕を組み合わせて最終動画を生成
"""
import asyncio
import shutil
import subprocess
import tempfile
from typing import Optional, Dict, Any, List
from pathlib import Path
from dataclasses import dataclass
import json
from datetime import datetime
from PIL import Image, ImageDraw

# 基本的なロガー設定（loguruの代替）
from core.utils.logger import logger
from core.utils.ffmpeg_utils import check_ffmpeg_with_warning, FFMPEG_INSTALL_GUIDE

from config.settings import settings
from notebook_lm.audio_generator import AudioInfo
from notebook_lm.transcript_processor import TranscriptInfo
from slides.slide_generator import SlidesPackage
from .subtitle_generator import SubtitleGenerator
from .effect_processor import EffectProcessor

@dataclass
class VideoInfo:
    """動画情報

    互換性のため、テストで使用されている `format` 引数をオプションフィールドとして受け付ける。
    既存コードでは解像度やコーデック情報を別途管理しているため、format はメタデータ用途とし、
    省略時は "mp4" をデフォルトとする。
    """

    file_path: Path
    duration: float
    resolution: tuple
    fps: int = 30
    file_size: int = 0
    has_subtitles: bool = False
    has_effects: bool = False
    created_at: datetime = datetime.now()
    format: str = "mp4"


@dataclass
class ThumbnailInfo:
    """サムネイル情報"""
    file_path: Path
    title_text: str
    subtitle_text: str
    style: str
    resolution: tuple
    file_size: int
    has_overlay: bool
    has_text_effects: bool
    created_at: datetime


class VideoComposer:
    """動画合成クラス"""
    
    def __init__(self):
        self.video_settings = settings.VIDEO_SETTINGS
        self.output_dir = settings.VIDEOS_DIR
        self.subtitle_generator = SubtitleGenerator()
        self.effect_processor = EffectProcessor()
        
    async def compose_video(
        self,
        audio_file: AudioInfo,
        slides_file: SlidesPackage,
        transcript: TranscriptInfo,
        quality: str = "1080p",
        timeline_plan: Optional[Dict[str, Any]] = None
    ) -> VideoInfo:
        """
        音声、スライド、字幕を合成して動画を生成
        
        Args:
            audio_file: 音声ファイル情報
            slides_file: スライドファイル情報
            transcript: 台本情報
            quality: 動画品質
            timeline_plan: タイムライン計画（セグメント長・エフェクト指示）
            
        Returns:
            VideoInfo: 生成された動画情報
        """
        logger.info("動画合成開始")
        
        # timeline_plan からセグメント情報を抽出
        self._current_timeline_plan = timeline_plan
        
        try:
            # Step 1: スライド画像を抽出
            slide_images = await self._extract_slide_images(slides_file)
            
            # Step 2: 字幕を生成
            subtitle_file = await self.subtitle_generator.generate_subtitles(transcript)
            
            # Step 3: スライドにエフェクトを適用
            processed_slides = await self.effect_processor.apply_effects(
                slide_images, transcript
            )
            
            # Step 4: 動画合成実行
            video_info = await self._compose_final_video(
                audio_file, processed_slides, subtitle_file, quality
            )
            
            # Step 5: メタデータ保存
            await self._save_video_metadata(video_info, transcript)
            
            logger.success(f"動画合成完了: {video_info.file_path}")
            return video_info
            
        except (OSError, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.error(f"動画合成エラー: {str(e)}")
            raise
        except Exception as e:
            logger.error(f"動画合成エラー: {str(e)}")
            raise
    
    async def _extract_slide_images(self, slides_file: SlidesPackage) -> List[Path]:
        """
        スライドファイルから画像を抽出
        
        Args:
            slides_file: スライドファイル情報
            
        Returns:
            List[Path]: 抽出された画像ファイルパス一覧
        """
        logger.info("スライド画像抽出中...")
        
        # 1) Google Slides のエクスポート画像を優先
        try:
            if getattr(slides_file, "presentation_id", ""):
                images_dir = settings.SLIDES_IMAGES_DIR / slides_file.presentation_id
                if images_dir.exists():
                    exported = sorted(images_dir.glob("*.png"))
                    if exported:
                        logger.info(f"エクスポート済みスライド画像を使用: {len(exported)}枚 from {images_dir}")
                        return exported
        except (OSError, AttributeError, TypeError, ValueError) as exc:
            logger.debug(f"エクスポート済みスライド画像の探索に失敗（PPTX抽出へフォールバック）: {exc}")
            pass
        except Exception as exc:
            logger.debug(f"エクスポート済みスライド画像の探索に失敗（PPTX抽出へフォールバック）: {exc}")
            pass
        
        # 2) PPTXからの抽出
        pptx_images = await self._extract_from_pptx(slides_file)
        if pptx_images:
            logger.info(f"PPTXからスライド画像を抽出: {len(pptx_images)}枚")
            return pptx_images
        
        # 3) フォールバック: プレースホルダー生成（実在するPNGを作成）
        slide_images: List[Path] = []
        target_width, target_height = self.video_settings["resolution"]
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # テーマ取得
        theme_name = getattr(settings, "PLACEHOLDER_THEME", "dark")
        themes = getattr(settings, "PLACEHOLDER_THEMES", {})
        theme = themes.get(theme_name, themes.get("dark", {
            "background": (20, 20, 25),
            "title_color": (235, 235, 235),
            "speaker_color": (180, 200, 255),
            "body_color": (200, 200, 200),
            "label_color": (120, 120, 130),
            "accent_color": (100, 150, 255),
        }))
        
        logger.info(f"プレースホルダーテーマ: {theme_name}")
        
        for i, slide in enumerate(slides_file.slides, 1):
            image_path = self.output_dir / f"slide_{i:03d}.png"
            
            if not image_path.exists():
                img = Image.new('RGB', (target_width, target_height), color=theme["background"])
                draw = ImageDraw.Draw(img)
                
                # アクセントライン（上部）
                accent_color = theme.get("accent_color", (100, 150, 255))
                draw.rectangle([(0, 0), (target_width, 4)], fill=accent_color)
                
                title_text = slide.title or f"Slide {i}"

                body_text = slide.content or ""
                subtitle_text = self._wrap_placeholder_text(body_text)

                speakers = getattr(slide, "speakers", None) or []
                show_speakers = bool(speakers) and settings.SLIDES_SETTINGS.get("show_speaker_on_placeholder", False)
                if show_speakers:
                    speaker_text = " / ".join(dict.fromkeys(str(s) for s in speakers if s))
                else:
                    speaker_text = ""

                total_slides = getattr(slides_file, "total_slides", None) or len(slides_file.slides)
                duration_seconds = float(slide.estimated_duration or 10.0)
                label_text = f"{i:02d}/{total_slides:02d}  •  {duration_seconds:.0f}s"

                x_margin = 40
                y = 60
                draw.text((x_margin, y), title_text, fill=theme["title_color"])
                y += 50

                if speaker_text:
                    draw.text((x_margin, y), speaker_text, fill=theme["speaker_color"])
                    y += 40

                if subtitle_text:
                    draw.text((x_margin, y), subtitle_text, fill=theme["body_color"])
                    y += 40

                draw.text((x_margin, target_height - 80), label_text, fill=theme["label_color"])
                
                # アクセントライン（下部）
                draw.rectangle([(0, target_height - 4), (target_width, target_height)], fill=accent_color)
                
                img.save(image_path, format='PNG')
            
            slide_images.append(image_path)
        
        logger.info(f"スライド画像抽出完了: {len(slide_images)}枚 (placeholder)")
        return slide_images
    
    def _wrap_placeholder_text(
        self,
        text: str,
        max_chars_per_line: int = 26,
        max_lines: int = 3,
    ) -> str:
        """プレースホルダ用にテキストを簡易折り返しする

        - 日本語/英語混在を前提に、文字数ベースで固定幅に分割
        - 最大行数を超える場合は末尾に "..." を付与して省略を示す
        """
        normalized = (text or "").strip()
        if not normalized:
            return ""

        lines = []
        index = 0
        length = len(normalized)

        while index < length and len(lines) < max_lines:
            end = min(index + max_chars_per_line, length)
            line = normalized[index:end]
            lines.append(line)
            index = end

        # まだテキストが残っている場合は末尾に省略記号を付ける
        if index < length and lines:
            last = lines[-1]
            if len(last) >= 3:
                lines[-1] = last[:-3] + "..."
            else:
                lines[-1] = last + "..."

        return "\n".join(lines)

    async def _extract_from_pptx(self, slides_file: SlidesPackage) -> List[Path]:
        """PPTXファイルからスライド画像を抽出

        python-pptx または LibreOffice を使用して PPTX から PNG を生成します。

        Args:
            slides_file: スライドファイル情報

        Returns:
            List[Path]: 抽出された画像ファイルパス一覧。失敗時は空リスト。
        """
        pptx_path = getattr(slides_file, "pptx_path", None)
        if not pptx_path:
            pptx_path = getattr(slides_file, "file_path", None)
        if not pptx_path:
            return []

        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            return []

        presentation_id = getattr(slides_file, "presentation_id", "") or pptx_path.stem
        output_dir = settings.SLIDES_IMAGES_DIR / presentation_id
        output_dir.mkdir(parents=True, exist_ok=True)

        # 既にエクスポート済みならそれを使用
        existing = sorted(output_dir.glob("*.png"))
        if existing:
            logger.info(f"既存のPPTXエクスポート画像を使用: {len(existing)}枚")
            return existing

        # 方法1: LibreOffice を使用 (推奨: 高品質)
        if shutil.which("soffice"):
            try:
                with tempfile.TemporaryDirectory() as tmp_dir:
                    # LibreOffice で PDF に変換
                    pdf_result = subprocess.run([
                        "soffice", "--headless", "--convert-to", "pdf",
                        "--outdir", tmp_dir, str(pptx_path)
                    ], capture_output=True, timeout=120)

                    if pdf_result.returncode == 0:
                        pdf_files = list(Path(tmp_dir).glob("*.pdf"))
                        if pdf_files:
                            pdf_path = pdf_files[0]
                            # PDF から PNG に変換 (pdftoppm / ImageMagick)
                            images = await self._pdf_to_images(pdf_path, output_dir)
                            if images:
                                logger.info(f"LibreOffice経由でPPTX抽出完了: {len(images)}枚")
                                return images
            except (subprocess.TimeoutExpired, OSError, ValueError, TypeError, RuntimeError) as e:
                logger.warning(f"LibreOfficeでのPPTX変換に失敗: {e}")
            except Exception as e:
                logger.warning(f"LibreOfficeでのPPTX変換に失敗: {e}")

        # 方法2: python-pptx を使用 (サムネイル抽出のみ)
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation(str(pptx_path))
            images: List[Path] = []

            for i, slide in enumerate(prs.slides, 1):
                # スライドから埋め込み画像を抽出
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_bytes = shape.image.blob
                            image_path = output_dir / f"slide_{i:03d}.png"
                            with open(image_path, "wb") as f:
                                f.write(image_bytes)
                            images.append(image_path)
                            break  # 1スライドにつき1画像
                        except (AttributeError, OSError, TypeError, ValueError) as exc:
                            logger.debug(f"python-pptx 画像抽出に失敗（スキップ）: {exc}")
                            continue

            if images:
                logger.info(f"python-pptxでPPTX抽出完了: {len(images)}枚")
                return images
        except ImportError:
            logger.debug("python-pptx がインストールされていません")
        except (OSError, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.warning(f"python-pptxでのPPTX抽出に失敗: {e}")
        except Exception as e:
            logger.warning(f"python-pptxでのPPTX抽出に失敗: {e}")

        return []

    async def _pdf_to_images(self, pdf_path: Path, output_dir: Path) -> List[Path]:
        """PDFをPNG画像に変換

        pdftoppm (poppler-utils) または ImageMagick を使用します。
        """
        images: List[Path] = []

        # pdftoppm を試行
        if shutil.which("pdftoppm"):
            try:
                prefix = output_dir / "slide"
                result = subprocess.run([
                    "pdftoppm", "-png", "-r", "150",
                    str(pdf_path), str(prefix)
                ], capture_output=True, timeout=120)
                if result.returncode == 0:
                    images = sorted(output_dir.glob("slide-*.png"))
                    # ファイル名を正規化
                    for i, img in enumerate(images, 1):
                        new_name = output_dir / f"slide_{i:03d}.png"
                        img.rename(new_name)
                    images = sorted(output_dir.glob("slide_*.png"))
            except (subprocess.TimeoutExpired, OSError, ValueError, TypeError, RuntimeError) as e:
                logger.warning(f"pdftoppmでのPDF変換に失敗: {e}")
            except Exception as e:
                logger.warning(f"pdftoppmでのPDF変換に失敗: {e}")

        # ImageMagick を試行
        if not images and shutil.which("convert"):
            try:
                result = subprocess.run([
                    "convert", "-density", "150",
                    str(pdf_path), str(output_dir / "slide_%03d.png")
                ], capture_output=True, timeout=120)
                if result.returncode == 0:
                    images = sorted(output_dir.glob("slide_*.png"))
            except (subprocess.TimeoutExpired, OSError, ValueError, TypeError, RuntimeError) as e:
                logger.warning(f"ImageMagickでのPDF変換に失敗: {e}")
            except Exception as e:
                logger.warning(f"ImageMagickでのPDF変換に失敗: {e}")

        return images
    
    async def _compose_final_video(
        self,
        audio_info: AudioInfo,
        slide_images: List[Path],
        subtitle_file: Path,
        quality: str
    ) -> VideoInfo:
        """
        最終動画を合成
        
        Args:
            audio_info: 音声情報
            slide_images: スライド画像一覧
            subtitle_file: 字幕ファイル
            quality: 動画品質
            
        Returns:
            VideoInfo: 生成された動画情報
        """
        logger.info("最終動画合成中...")
        
        # 出力ファイルパス
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = self.output_dir / f"generated_video_{timestamp}.mp4"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # 動画設定
        resolution = self._get_resolution_from_quality(quality)
        fps = self.video_settings["fps"]
        
        try:
            # MoviePyを使用した動画合成
            from moviepy.editor import (
                AudioFileClip, ImageClip, CompositeVideoClip, 
                concatenate_videoclips, TextClip
            )
            
            # 音声読み込み
            audio_clip = AudioFileClip(str(audio_info.file_path))
            
            # スライド動画クリップ作成
            video_clips = []
            num_items = max(len(slide_images), 1)
            default_slide_duration = max(audio_clip.duration / num_items, 0.1)
            
            for i, slide_image in enumerate(slide_images):
                # ProcessedSlide にも対応（最初のフレームを使用）
                img_path = slide_image
                clip_duration = default_slide_duration
                try:
                    # dataclass で processed_frames を持つ場合
                    if hasattr(slide_image, "processed_frames") and getattr(slide_image, "processed_frames"):
                        img_path = slide_image.processed_frames[0]
                    
                    # 個別スライドのdurationを取得
                    if hasattr(slide_image, "duration"):
                        try:
                            duration_value = float(getattr(slide_image, "duration") or 0.0)
                            if duration_value > 0:
                                clip_duration = duration_value
                        except (TypeError, ValueError):
                            pass
                except (AttributeError, TypeError, ValueError) as exc:
                    logger.debug(f"スライド情報の読み取りに失敗（フォールバック）: {exc}")
                    img_path = slide_image
                except Exception as exc:
                    logger.debug(f"スライド情報の読み取りで予期しない例外（フォールバック）: {exc}")
                    img_path = slide_image
                
                # 画像クリップ作成
                img_clip = ImageClip(str(img_path))
                img_clip = img_clip.set_duration(clip_duration)
                img_clip = img_clip.resize(resolution)
                
                video_clips.append(img_clip)
            
            video_clip = concatenate_videoclips(video_clips)
            
            # 音声を設定
            final_clip = video_clip.set_audio(audio_clip)
            
            # 字幕を追加
            if subtitle_file.exists():
                final_clip = self._add_subtitles_to_video(final_clip, subtitle_file)
            
            # 動画出力
            final_clip.write_videofile(
                str(output_path),
                fps=fps,
                codec=self.video_settings["video_codec"],
                audio_codec=self.video_settings["audio_codec"],
                temp_audiofile="temp-audio.m4a",
                remove_temp=True
            )
            
            # リソース解放
            audio_clip.close()
            final_clip.close()
            
            # 動画情報作成
            video_info = VideoInfo(
                file_path=output_path,
                duration=audio_clip.duration,
                resolution=resolution,
                fps=fps,
                file_size=output_path.stat().st_size,
                has_subtitles=subtitle_file.exists(),
                has_effects=True,
                created_at=datetime.now()
            )
            
            logger.info("最終動画合成完了")
            return video_info
            
        except ImportError:
            logger.error("MoviePyライブラリが見つかりません")
            # フォールバック実装
            return await self._compose_video_fallback(
                audio_info, slide_images, subtitle_file, quality, output_path
            )
        except (OSError, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.warning(f"MoviePy処理で例外が発生したためフォールバックします: {e}")
            return await self._compose_video_fallback(
                audio_info, slide_images, subtitle_file, quality, output_path
            )
        except Exception as e:
            logger.warning(f"MoviePy処理で例外が発生したためフォールバックします: {e}")
            return await self._compose_video_fallback(
                audio_info, slide_images, subtitle_file, quality, output_path
            )
    
    def _get_resolution_from_quality(self, quality: str) -> tuple:
        """
        品質設定から解像度を取得
        
        Args:
            quality: 品質設定
            
        Returns:
            tuple: 解像度 (width, height)
        """
        quality_map = {
            "720p": (1280, 720),
            "1080p": (1920, 1080),
            "4k": (3840, 2160)
        }
        return quality_map.get(quality, (1920, 1080))
    
    def _add_subtitles_to_video(self, video_clip, subtitle_file: Path):
        """動画に字幕をオーバーレイする

        - `pysrt` がインストールされている場合: MoviePy + TextClip でハードサブタイトルを合成
        - `pysrt` が無い場合: 処理をスキップし、外部字幕ファイル(SRT/ASS/VTT)のみを利用する
        """
        try:
            from moviepy.editor import CompositeVideoClip, TextClip

            try:
                import pysrt
            except ImportError:
                # 環境に pysrt が無い場合は動画への埋め込みを行わず、外部字幕としての利用に限定する
                logger.info(
                    "pysrt がインストールされていないため、動画への字幕オーバーレイをスキップします。"
                    "SRT/ASS/VTT ファイルは外部字幕として再生ソフト側で利用できます。"
                )
                return video_clip

            # SRTファイル読み込み
            subs = pysrt.open(str(subtitle_file))

            subtitle_clips = []
            for sub in subs:
                # 時間変換
                start_time = self._srt_time_to_seconds(sub.start)
                end_time = self._srt_time_to_seconds(sub.end)

                # テキストクリップ作成
                txt_clip = TextClip(
                    sub.text,
                    fontsize=settings.SUBTITLE_SETTINGS["font_size"],
                    color=settings.SUBTITLE_SETTINGS["font_color"],
                    font=settings.SUBTITLE_SETTINGS["font_family"],
                ).set_position(("center", "bottom")).set_start(start_time).set_end(end_time)

                subtitle_clips.append(txt_clip)

            # 字幕を動画に合成
            return CompositeVideoClip([video_clip] + subtitle_clips)

        except (ImportError, OSError, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.warning(f"字幕追加に失敗: {str(e)}")
            return video_clip
        except Exception as e:
            logger.warning(f"字幕追加に失敗: {str(e)}")
            return video_clip
    
    def _srt_time_to_seconds(self, srt_time) -> float:
        """
        SRT時間を秒に変換
        
        Args:
            srt_time: SRT時間オブジェクト
            
        Returns:
            float: 秒数
        """
        return (srt_time.hours * 3600 + 
                srt_time.minutes * 60 + 
                srt_time.seconds + 
                srt_time.milliseconds / 1000.0)
    
    async def _compose_video_fallback(
        self,
        audio_info: AudioInfo,
        slide_images: List[Path],
        subtitle_file: Path,
        quality: str,
        output_path: Path
    ) -> VideoInfo:
        """
        フォールバック動画合成（FFmpegを使用）
        
        Args:
            audio_info: 音声情報
            slide_images: スライド画像一覧
            subtitle_file: 字幕ファイル
            quality: 動画品質
            output_path: 出力パス
            
        Returns:
            VideoInfo: 生成された動画情報
        """
        logger.info("FFmpegを使用したフォールバック動画合成")
        
        resolution = self._get_resolution_from_quality(quality)
        fps = self.video_settings.get("fps", 30)
        
        # FFmpegが利用可能か確認（改善されたユーティリティを使用）
        ffmpeg_available, ffmpeg_path = check_ffmpeg_with_warning()
        if not ffmpeg_available:
            logger.error("FFmpegが見つかりません。動画出力をスキップし、メタデータのみ保存します。")
            logger.error("動画出力を有効にするには、上記のインストール手順に従ってFFmpegをインストールしてください。")
            with open(output_path, 'wb') as f:
                f.write(b'')
            return VideoInfo(
                file_path=output_path,
                duration=audio_info.duration,
                resolution=resolution,
                fps=fps,
                file_size=0,
                has_subtitles=False,
                has_effects=False,
                created_at=datetime.now()
            )
        
        try:
            with tempfile.TemporaryDirectory() as tmp_dir:
                tmp_path = Path(tmp_dir)
                
                # スライド画像からビデオを作成
                num_slides = max(len(slide_images), 1)
                slide_duration = max(audio_info.duration / num_slides, 0.1)
                
                # concat用のファイルリストを作成
                concat_file = tmp_path / "concat.txt"
                with open(concat_file, 'w', encoding='utf-8') as f:
                    for slide_image in slide_images:
                        # パスから画像を取得
                        img_path = slide_image
                        if hasattr(slide_image, "processed_frames") and getattr(slide_image, "processed_frames"):
                            img_path = slide_image.processed_frames[0]
                        
                        # 個別スライドのdurationを取得
                        dur = slide_duration
                        if hasattr(slide_image, "duration"):
                            try:
                                dur = float(getattr(slide_image, "duration") or slide_duration)
                            except (TypeError, ValueError):
                                pass
                        
                        f.write(f"file '{img_path}'\n")
                        f.write(f"duration {dur}\n")
                    # 最後のフレームを保持するため再度追加
                    if slide_images:
                        last_img = slide_images[-1]
                        if hasattr(last_img, "processed_frames") and getattr(last_img, "processed_frames"):
                            last_img = last_img.processed_frames[0]
                        f.write(f"file '{last_img}'\n")
                
                # 一時動画ファイル
                temp_video = tmp_path / "temp_video.mp4"
                
                # Step 1: 画像から動画を作成
                cmd_video = [
                    ffmpeg_path, "-y",
                    "-f", "concat", "-safe", "0",
                    "-i", str(concat_file),
                    "-vf", f"scale={resolution[0]}:{resolution[1]}:force_original_aspect_ratio=decrease,pad={resolution[0]}:{resolution[1]}:(ow-iw)/2:(oh-ih)/2",
                    "-c:v", "libx264", "-preset", "medium",
                    "-pix_fmt", "yuv420p",
                    "-r", str(fps),
                    str(temp_video)
                ]
                
                result = subprocess.run(cmd_video, capture_output=True, timeout=300)
                if result.returncode != 0:
                    logger.warning(f"FFmpeg動画生成エラー: {result.stderr.decode()}")
                    raise RuntimeError("FFmpeg video generation failed")
                
                # Step 2: 音声を追加
                temp_with_audio = tmp_path / "temp_with_audio.mp4"
                cmd_audio = [
                    ffmpeg_path, "-y",
                    "-i", str(temp_video),
                    "-i", str(audio_info.file_path),
                    "-c:v", "copy",
                    "-c:a", "aac", "-b:a", "192k",
                    "-shortest",
                    str(temp_with_audio)
                ]
                
                result = subprocess.run(cmd_audio, capture_output=True, timeout=300)
                if result.returncode != 0:
                    logger.warning(f"FFmpeg音声追加エラー: {result.stderr.decode()}")
                    raise RuntimeError("FFmpeg audio addition failed")
                
                # Step 3: 字幕を追加（オプション）
                has_subtitles = False
                final_temp = temp_with_audio
                
                if subtitle_file.exists() and subtitle_file.stat().st_size > 0:
                    temp_with_subs = tmp_path / "temp_with_subs.mp4"
                    # 字幕をハードコード（焼き込み）
                    cmd_subs = [
                        ffmpeg_path, "-y",
                        "-i", str(temp_with_audio),
                        "-vf", f"subtitles='{str(subtitle_file).replace(chr(92), chr(92)+chr(92)).replace(':', chr(92)+':')}'",
                        "-c:a", "copy",
                        str(temp_with_subs)
                    ]
                    
                    result = subprocess.run(cmd_subs, capture_output=True, timeout=300)
                    if result.returncode == 0:
                        final_temp = temp_with_subs
                        has_subtitles = True
                        logger.info("字幕を動画に焼き込みました")
                    else:
                        logger.warning(f"字幕追加をスキップ: {result.stderr.decode()[:200]}")
                
                # 最終出力にコピー
                shutil.copy2(final_temp, output_path)
                
                file_size = output_path.stat().st_size if output_path.exists() else 0
                logger.info(f"FFmpegによる動画合成完了: {output_path} ({file_size} bytes)")
                
                return VideoInfo(
                    file_path=output_path,
                    duration=audio_info.duration,
                    resolution=resolution,
                    fps=fps,
                    file_size=file_size,
                    has_subtitles=has_subtitles,
                    has_effects=False,
                    created_at=datetime.now()
                )
                
        except (subprocess.TimeoutExpired, OSError, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.error(f"FFmpegフォールバック動画合成に失敗: {e}")
            # 最終フォールバック: 空のファイル
            with open(output_path, 'wb') as f:
                f.write(b'')
            return VideoInfo(
                file_path=output_path,
                duration=audio_info.duration,
                resolution=resolution,
                fps=fps,
                file_size=0,
                has_subtitles=False,
                has_effects=False,
                created_at=datetime.now()
            )
        except Exception as e:
            logger.error(f"FFmpegフォールバック動画合成に失敗: {e}")
            # 最終フォールバック: 空のファイル
            with open(output_path, 'wb') as f:
                f.write(b'')
            return VideoInfo(
                file_path=output_path,
                duration=audio_info.duration,
                resolution=resolution,
                fps=fps,
                file_size=0,
                has_subtitles=False,
                has_effects=False,
                created_at=datetime.now()
            )
    
    async def _save_video_metadata(self, video_info: VideoInfo, transcript: TranscriptInfo):
        """
        動画メタデータを保存
        
        Args:
            video_info: 動画情報
            transcript: 台本情報
        """
        metadata_path = video_info.file_path.with_suffix('.json')
        
        metadata = {
            "video_file": str(video_info.file_path),
            "duration": video_info.duration,
            "resolution": video_info.resolution,
            "fps": video_info.fps,
            "file_size": video_info.file_size,
            "has_subtitles": video_info.has_subtitles,
            "has_effects": video_info.has_effects,
            "created_at": video_info.created_at.isoformat(),
            "source_transcript": transcript.title,
            "total_segments": len(transcript.segments),
            "transcript_accuracy": transcript.accuracy_score
        }
        
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
        
        logger.info(f"動画メタデータ保存完了: {metadata_path}")
    
    def optimize_for_platform(self, video_info: VideoInfo, platform: str = "youtube") -> VideoInfo:
        """
        プラットフォーム向けに動画を最適化
        
        Args:
            video_info: 動画情報
            platform: 対象プラットフォーム
            
        Returns:
            VideoInfo: 最適化された動画情報
        """
        logger.info(f"{platform}向け最適化実行")
        
        # プラットフォーム別最適化設定
        platform_settings = {
            "youtube": {
                "max_file_size": 128 * 1024 * 1024 * 1024,  # 128GB
                "recommended_bitrate": "8000k",
                "recommended_fps": 30
            },
            "tiktok": {
                "max_file_size": 4 * 1024 * 1024 * 1024,  # 4GB
                "recommended_aspect_ratio": (9, 16),
                "max_duration": 600  # 10分
            }
        }
        
        settings = platform_settings.get(platform, platform_settings["youtube"])
        
        # ファイルサイズチェック
        if video_info.file_size > settings["max_file_size"]:
            logger.warning(f"ファイルサイズが制限を超えています: {video_info.file_size / (1024*1024):.1f}MB > {settings['max_file_size'] / (1024*1024):.1f}MB")
            
            # 圧縮処理
            compressed_path = self._compress_video(
                video_info.video_path,
                target_size=settings["max_file_size"],
                platform=platform
            )
            
            if compressed_path and compressed_path.exists():
                video_info.video_path = compressed_path
                video_info.file_size = compressed_path.stat().st_size
                logger.info(f"圧縮完了: {video_info.file_size / (1024*1024):.1f}MB")
        
        return video_info
    
    def _compress_video(
        self,
        video_path: Path,
        target_size: int,
        platform: str = "youtube"
    ) -> Optional[Path]:
        """動画を圧縮
        
        Args:
            video_path: 元動画パス
            target_size: 目標ファイルサイズ（バイト）
            platform: プラットフォーム
            
        Returns:
            圧縮後の動画パス（失敗時はNone）
        """
        try:
            output_path = video_path.parent / f"{video_path.stem}_compressed{video_path.suffix}"
            
            # 現在のファイルサイズと目標サイズから圧縮率を計算
            current_size = video_path.stat().st_size
            target_bitrate = int((target_size * 8) / self._get_video_duration(video_path) * 0.9)  # 90%マージン
            
            # FFmpegで圧縮
            cmd = [
                "ffmpeg", "-y",
                "-i", str(video_path),
                "-c:v", "libx264",
                "-preset", "medium",
                "-b:v", f"{target_bitrate}",
                "-c:a", "aac",
                "-b:a", "128k",
                "-movflags", "+faststart",
                str(output_path)
            ]
            
            logger.info(f"動画圧縮開始: 目標ビットレート {target_bitrate}bps")
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode == 0 and output_path.exists():
                # 圧縮後もサイズオーバーならCRF方式で再圧縮
                if output_path.stat().st_size > target_size:
                    logger.warning("ビットレート指定では目標サイズに収まらず、CRF方式で再圧縮")
                    output_path.unlink()
                    
                    cmd_crf = [
                        "ffmpeg", "-y",
                        "-i", str(video_path),
                        "-c:v", "libx264",
                        "-preset", "slow",
                        "-crf", "28",  # 高圧縮
                        "-c:a", "aac",
                        "-b:a", "96k",
                        "-movflags", "+faststart",
                        str(output_path)
                    ]
                    
                    result = subprocess.run(cmd_crf, capture_output=True, text=True)
                    
                    if result.returncode != 0 or not output_path.exists():
                        logger.error(f"CRF圧縮失敗: {result.stderr}")
                        return None
                
                logger.info(f"圧縮完了: {output_path}")
                return output_path
            else:
                logger.error(f"圧縮失敗: {result.stderr}")
                return None
                
        except FileNotFoundError:
            logger.warning("FFmpegが見つかりません。圧縮をスキップします。")
            logger.warning("FFmpegをインストールするには、以下を参照してください:")
            logger.warning("  Windows: winget install FFmpeg")
            logger.warning("  macOS: brew install ffmpeg")
            logger.warning("  Linux: sudo apt install ffmpeg")
            return None
        except (OSError, subprocess.TimeoutExpired, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.error(f"圧縮エラー: {e}")
            return None
        except Exception as e:
            logger.error(f"圧縮エラー: {e}")
            return None
    
    def _get_video_duration(self, video_path: Path) -> float:
        """動画の長さを取得（秒）"""
        try:
            cmd = [
                "ffprobe",
                "-v", "error",
                "-show_entries", "format=duration",
                "-of", "default=noprint_wrappers=1:nokey=1",
                str(video_path)
            ]
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode == 0:
                return float(result.stdout.strip())
        except (OSError, subprocess.TimeoutExpired, AttributeError, TypeError, ValueError, RuntimeError) as exc:
            logger.debug(f"ffprobe による動画長取得に失敗（フォールバック=600秒）: {exc}")
            pass
        except Exception as exc:
            logger.debug(f"ffprobe による動画長取得に失敗（フォールバック=600秒）: {exc}")
            pass
        
        # フォールバック: 10分と仮定
        return 600.0

    async def generate_thumbnail(
        self,
        title: str,
        first_slide_path: Optional[Path] = None,
        output_path: Optional[Path] = None,
    ) -> Optional[Path]:
        """YouTube用サムネイルを生成
        
        Args:
            title: 動画タイトル
            first_slide_path: 最初のスライド画像パス（省略時は新規生成）
            output_path: 出力パス（省略時は自動決定）
            
        Returns:
            サムネイル画像パス（失敗時はNone）
        """
        logger.info(f"サムネイル生成開始: {title}")
        
        try:
            # 出力パス決定
            if output_path is None:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = settings.THUMBNAILS_DIR / f"thumbnail_{timestamp}.png"
            
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            # YouTube推奨サイズ: 1280x720
            thumb_width, thumb_height = 1280, 720
            
            # ベース画像を作成または読み込み
            if first_slide_path and first_slide_path.exists():
                base_img = Image.open(first_slide_path)
                base_img = base_img.resize((thumb_width, thumb_height), Image.Resampling.LANCZOS)
            else:
                # グラデーション背景を生成
                base_img = Image.new('RGB', (thumb_width, thumb_height), color=(30, 30, 40))
                draw = ImageDraw.Draw(base_img)
                
                # シンプルなグラデーション効果
                for y in range(thumb_height):
                    alpha = int(y / thumb_height * 60)
                    draw.line([(0, y), (thumb_width, y)], fill=(30 + alpha, 30 + alpha, 50 + alpha))
            
            draw = ImageDraw.Draw(base_img)
            
            # タイトル用の半透明オーバーレイ（下部）
            overlay_height = 180
            overlay = Image.new('RGBA', (thumb_width, overlay_height), (0, 0, 0, 180))
            base_img.paste(
                Image.alpha_composite(
                    base_img.crop((0, thumb_height - overlay_height, thumb_width, thumb_height)).convert('RGBA'),
                    overlay
                ).convert('RGB'),
                (0, thumb_height - overlay_height)
            )
            
            # タイトルテキストを描画
            draw = ImageDraw.Draw(base_img)
            
            # テキストを折り返し
            max_chars_per_line = 20
            lines = []
            words = title
            while words:
                if len(words) <= max_chars_per_line:
                    lines.append(words)
                    break
                else:
                    lines.append(words[:max_chars_per_line])
                    words = words[max_chars_per_line:]
            
            # テキスト描画（中央揃え）
            text_y = thumb_height - overlay_height + 30
            for line in lines[:3]:  # 最大3行
                text_width = len(line) * 24  # 概算
                text_x = (thumb_width - text_width) // 2
                draw.text((text_x, text_y), line, fill=(255, 255, 255))
                text_y += 50
            
            # 保存
            base_img.save(output_path, format='PNG', quality=95)
            logger.success(f"サムネイル生成完了: {output_path}")
            
            return output_path
            
        except (OSError, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.error(f"サムネイル生成エラー: {e}")
            return None
        except Exception as e:
            logger.error(f"サムネイル生成エラー: {e}")
            return None
