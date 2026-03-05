"""
スライド生成モジュール
Google Slidesを使用したスライドの自動生成
"""
from typing import List, Dict, Any, Optional
from pathlib import Path
from dataclasses import dataclass
import time
import json

# 基本的なロガー設定（loguruの代替）
from core.utils.logger import logger

from config.settings import settings
from notebook_lm.transcript_processor import TranscriptInfo
from .content_splitter import ContentSplitter
from .google_slides_client import GoogleSlidesClient

@dataclass
class SlideInfo:
    """スライド情報"""
    slide_id: int
    title: str
    content: str
    layout_type: Optional[str] = None
    estimated_duration: Optional[float] = None
    image_suggestions: List[str] = None
    # 互換性維持のためのフィールド（旧API: layout/duration）
    layout: str = None
    duration: float = None
    speakers: Optional[List[str]] = None

    def __post_init__(self):
        # 旧フィールドが渡された場合、新フィールドに反映
        if self.layout and not self.layout_type:
            self.layout_type = self.layout
        if self.duration and not self.estimated_duration:
            self.estimated_duration = self.duration

@dataclass
class SlidesPackage:
    """スライドパッケージ情報

    テスト互換性のため、file_path / total_slides / theme は省略可能な引数として扱う。
    既存のテストでは presentation_id と slides のみを指定しているため、
    それら以外のフィールドにはデフォルト値を設定する。
    """

    file_path: Optional[Path] = None
    slides: List[SlideInfo] = None
    total_slides: int = 0
    theme: str = "default"
    presentation_id: str = ""
    title: str = ""
    created_at: str = None

    def __post_init__(self):
        if self.slides is None:
            self.slides = []

class SlideGenerator:
    """スライド生成クラス"""
    
    def __init__(self):
        self.max_chars_per_slide = settings.SLIDES_SETTINGS["max_chars_per_slide"]
        self.max_slides_per_batch = settings.SLIDES_SETTINGS["max_slides_per_batch"]
        self.theme = settings.SLIDES_SETTINGS["theme"]
        self.output_dir = settings.SLIDES_DIR
        self.content_splitter = ContentSplitter()
    
    async def authenticate(self) -> bool:
        """Google Slides API 認証"""
        client = GoogleSlidesClient()
        available = client.is_available()
        if available:
            logger.info("Google Slides API 認証済み")
        else:
            logger.warning("Google Slides API 未認証のため、モックへフォールバックします")
        return available
        
    async def generate_slides(
        self,
        transcript: TranscriptInfo,
        max_slides: int = 20,
        script_bundle: Optional[Dict[str, Any]] = None
    ) -> SlidesPackage:
        """
        台本からスライドを生成
        
        Args:
            transcript: 台本情報
            max_slides: 最大スライド数
            script_bundle: オプションのスクリプトバンドル（NotebookLM対応）
            
        Returns:
            SlidesPackage: 生成されたスライドパッケージ
        """
        logger.info(f"スライド生成開始: {transcript.title}")
        
        prefer_bundle = settings.SLIDES_SETTINGS.get("prefer_gemini_slide_content", False)
        has_bundle = script_bundle is not None
        has_slides_in_bundle = has_bundle and "slides" in script_bundle
        bundle_slide_count = len(script_bundle.get("slides", [])) if has_slides_in_bundle else 0
        
        logger.info(
            f"スライド生成パラメータ: prefer_bundle={prefer_bundle}, "
            f"has_bundle={has_bundle}, has_slides_in_bundle={has_slides_in_bundle}, "
            f"bundle_slide_count={bundle_slide_count}"
        )

        # NotebookLM / Gemini 由来スライド情報を優先使用する場合
        if prefer_bundle and script_bundle and "slides" in script_bundle:
            logger.info(f"NotebookLM/Gemini のスライド情報を優先使用します ({bundle_slide_count}枚)")
            return await self._generate_slides_from_bundle(script_bundle, max_slides)
        
        try:
            # Step 1: 台本をスライド用に分割
            slide_contents = await self.content_splitter.split_for_slides(
                transcript, max_slides
            )
            
            # Step 2: Google Slidesでスライド生成
            slides_package = await self._generate_slides_with_google(
                slide_contents, transcript.title
            )
            
            # Step 3: スライドファイルのダウンロード
            await self._download_slides_file(slides_package)
            
            # Step 4: スライド情報の保存
            await self._save_slides_metadata(slides_package)
            
            logger.success(f"スライド生成完了: {slides_package.total_slides}枚")
            return slides_package
            
        except (OSError, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.error(f"スライド生成エラー: {str(e)}")
            raise
        except Exception as e:
            logger.error(f"スライド生成エラー: {str(e)}")
            raise
    
    async def _generate_slides_with_google(
        self, 
        slide_contents: List[Dict[str, Any]], 
        presentation_title: str
    ) -> SlidesPackage:
        """
        Google Slidesでスライドを生成
        
        Args:
            slide_contents: スライド内容一覧
            presentation_title: プレゼンテーションタイトル
            
        Returns:
            SlidesPackage: 生成されたスライドパッケージ
        """
        logger.info("Google Slidesでスライド生成中...")
        client = GoogleSlidesClient()
        presentation_id: Optional[str] = None
        slides: List[SlideInfo] = []
        
        # まずはAPIが利用可能か試みる
        try:
            presentation_id = client.create_presentation(presentation_title)
        except (AttributeError, TypeError, ValueError, OSError, RuntimeError) as e:
            presentation_id = None
            logger.warning(f"Slides APIのプレゼン作成に失敗: {e}")
        except Exception as e:
            presentation_id = None
            logger.warning(f"Slides APIのプレゼン作成に失敗: {e}")
        
        if presentation_id:
            # API経由でのスライド追加（簡易）
            try:
                simplified = [
                    {
                        "title": c.get("title", f"スライド {i+1}"),
                        "content": c.get("text", ""),
                        "duration": c.get("duration", 15.0)
                    }
                    for i, c in enumerate(slide_contents)
                ]
                client.add_slides(presentation_id, simplified)
            except (AttributeError, TypeError, ValueError, OSError, RuntimeError) as e:
                logger.warning(f"スライド追加でエラー（フォールバック継続）: {e}")
            except Exception as e:
                logger.warning(f"スライド追加でエラー（フォールバック継続）: {e}")
            
            # SlideInfoを整備
            for i, c in enumerate(slide_contents, start=1):
                slides.append(SlideInfo(
                    slide_id=i,
                    title=c.get("title", f"スライド {i}"),
                    content=c.get("text", ""),
                    layout_type=c.get("layout", "title_and_content"),
                    estimated_duration=c.get("duration", 15.0),
                    speakers=c.get("speakers"),
                ))
            
            # PPTXとサムネイル画像を書き出し
            pptx_path = self.output_dir / f"{presentation_id}.pptx"
            try:
                client.export_pptx(presentation_id, pptx_path)
            except (ImportError, AttributeError, TypeError, ValueError, OSError, RuntimeError) as e:
                logger.warning(f"PPTXエクスポート失敗: {e}")
            except Exception as e:
                logger.warning(f"PPTXエクスポート失敗: {e}")
            
            try:
                client.export_thumbnails(presentation_id, settings.SLIDES_IMAGES_DIR / presentation_id)
            except (ImportError, AttributeError, TypeError, ValueError, OSError, RuntimeError) as e:
                logger.warning(f"サムネイル書き出し失敗: {e}")
            except Exception as e:
                logger.warning(f"サムネイル書き出し失敗: {e}")
            
            slides_package = SlidesPackage(
                file_path=pptx_path,
                slides=slides,
                total_slides=len(slides),
                theme=self.theme,
                presentation_id=presentation_id,
                title=presentation_title,
                created_at=time.strftime("%Y-%m-%d %H:%M:%S")
            )
            logger.info(f"Google Slidesでの生成完了: {len(slides)}枚")
            return slides_package
        
        # APIが使えない場合は既存のモック実装にフォールバック
        logger.warning("Slides API未使用のため、モック生成にフォールバックします")
        slides = [
            SlideInfo(
                slide_id=i,
                title=c.get("title", f"スライド {i}"),
                content=c.get("text", c.get("content", "")),
                layout_type=c.get("layout", "title_and_content"),
                estimated_duration=c.get("duration", 15.0),
                speakers=c.get("speakers"),
            )
            for i, c in enumerate(slide_contents, start=1)
        ]
        presentation_id = f"mock_{int(time.time())}"
        slides_package = SlidesPackage(
            file_path=self.output_dir / f"{presentation_id}.pptx",
            slides=slides,
            total_slides=len(slides),
            theme=self.theme,
            presentation_id=presentation_id,
            title=presentation_title,
            created_at=time.strftime("%Y-%m-%d %H:%M:%S")
        )
        return slides_package
    
    async def _generate_slides_from_bundle(
        self,
        script_bundle: Dict[str, Any],
        max_slides: int
    ) -> SlidesPackage:
        """
        スクリプトバンドルからスライドを生成（NotebookLM対応）
        
        Args:
            script_bundle: スクリプトバンドル
            max_slides: 最大スライド数
            
        Returns:
            SlidesPackage: 生成されたスライドパッケージ
        """
        logger.info("スクリプトバンドルからスライド生成中...")
        
        bundle_slides = script_bundle.get("slides", [])
        slides: List[SlideInfo] = []
        
        # バンドルのスライドを SlideInfo に変換
        for i, bundle_slide in enumerate(bundle_slides[:max_slides]):
            slide_info = SlideInfo(
                slide_id=i + 1,
                title=bundle_slide.get("title", f"スライド {i+1}"),
                content=bundle_slide.get("content", ""),
                layout_type=bundle_slide.get("layout", "title_and_content"),
                estimated_duration=bundle_slide.get("duration", 15.0),
                # Gemini/NotebookLM 由来のフィールド名差異を吸収
                image_suggestions=bundle_slide.get(
                    "image_suggestions", bundle_slide.get("images", [])
                ),
            )
            slides.append(slide_info)
        
        # スライドが足りない場合はセグメントから補完
        if len(slides) < max_slides and "segments" in script_bundle:
            segments = script_bundle["segments"]
            for j, segment in enumerate(segments[len(slides):max_slides]):
                slide_info = SlideInfo(
                    slide_id=len(slides) + j + 1,
                    title=f"セグメント {len(slides) + j + 1}",
                    content=segment.get("content", ""),
                    layout_type="title_and_content",
                    estimated_duration=segment.get("duration", 15.0)
                )
                slides.append(slide_info)
        
        presentation_id = f"notebooklm_{int(time.time())}"
        title = script_bundle.get("title", "NotebookLM Presentation")
        
        slides_package = SlidesPackage(
            file_path=self.output_dir / f"{presentation_id}.pptx",
            slides=slides,
            total_slides=len(slides),
            theme=self.theme,
            presentation_id=presentation_id,
            title=title,
            created_at=time.strftime("%Y-%m-%d %H:%M:%S")
        )
        
        # ファイル生成とメタデータ保存
        await self._download_slides_file(slides_package)
        await self._save_slides_metadata(slides_package)
        
        logger.info(f"バンドルからのスライド生成完了: {len(slides)}枚")
        return slides_package

    async def create_slides_from_content(
        self,
        slides_content: List[Dict[str, Any]],
        presentation_title: str
    ) -> SlidesPackage:
        """テスト/CSVタイムライン用のシンプルなスライド生成"""
        slides: List[SlideInfo] = []
        for i, content in enumerate(slides_content, start=1):
            slide = SlideInfo(
                slide_id=content.get("slide_id", i),
                title=content.get("title", f"スライド {i}"),
                content=content.get("content", content.get("text", "")),
                layout_type=content.get("layout_type", content.get("layout")),
                estimated_duration=content.get("duration", 15.0),
                image_suggestions=content.get("image_suggestions", content.get("images")),
                speakers=content.get("speakers"),
            )
            slides.append(slide)

        self.output_dir.mkdir(parents=True, exist_ok=True)
        presentation_id = f"presentation_{int(time.time())}"
        slides_package = SlidesPackage(
            file_path=self.output_dir / f"{presentation_id}.pptx",
            slides=slides,
            total_slides=len(slides),
            theme=self.theme,
            presentation_id=presentation_id,
            title=presentation_title,
            created_at=time.strftime("%Y-%m-%d %H:%M:%S"),
        )
        await self._download_slides_file(slides_package)
        await self._save_slides_metadata(slides_package)
        return slides_package
    
    async def _download_slides_file(self, slides_package: SlidesPackage):
        """
        スライドファイルをダウンロード
        
        Args:
            slides_package: スライドパッケージ
        """
        logger.info("スライドファイルダウンロード中...")
        
        # 既にAPIでエクスポート済みなら何もしない
        if slides_package.file_path.exists():
            logger.info(f"既存スライドを検出: {slides_package.file_path} -> ダウンロード処理をスキップ")
            return
        
        # ディレクトリ作成
        slides_package.file_path.parent.mkdir(parents=True, exist_ok=True)
 
        try:
            from pptx import Presentation
        except ImportError as e:
            logger.warning(f"python-pptx が見つからないため PPTX 生成をスキップします: {e}")
            with open(slides_package.file_path, "wb") as f:
                f.write(b"")
            logger.info(f"スライドダウンロード完了: {slides_package.file_path}")
            return
 
        try:
            prs = Presentation()
            slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            for slide_info in slides_package.slides:
                slide = prs.slides.add_slide(slide_layout)
                try:
                    if getattr(slide.shapes, "title", None) is not None:
                        slide.shapes.title.text = slide_info.title or ""
                except (AttributeError, TypeError, ValueError):
                    pass
                try:
                    placeholders = getattr(slide.shapes, "placeholders", None)
                    if placeholders and len(placeholders) > 1:
                        placeholders[1].text = slide_info.content or ""
                except (AttributeError, TypeError, ValueError, IndexError, KeyError):
                    pass
 
            prs.save(str(slides_package.file_path))
            logger.info(f"スライドダウンロード完了: {slides_package.file_path}")
        except (OSError, AttributeError, TypeError, ValueError, RuntimeError) as e:
            logger.warning(f"PPTX生成に失敗しました: {e}")
            with open(slides_package.file_path, "wb") as f:
                f.write(b"")
            logger.info(f"スライドダウンロード完了: {slides_package.file_path}")
        except Exception as e:
            logger.warning(f"PPTX生成に失敗しました: {e}")
            with open(slides_package.file_path, "wb") as f:
                f.write(b"")
            logger.info(f"スライドダウンロード完了: {slides_package.file_path}")
    
    async def _save_slides_metadata(self, slides_package: SlidesPackage):
        """
        スライドメタデータを保存
        
        Args:
            slides_package: スライドパッケージ
        """
        metadata_path = self.output_dir / f"{slides_package.presentation_id}_metadata.json"
        
        metadata = {
            "presentation_id": slides_package.presentation_id,
            "title": slides_package.title,
            "total_slides": slides_package.total_slides,
            "created_at": slides_package.created_at,
            "file_path": str(slides_package.file_path),
            "slides": [
                {
                    "slide_id": slide.slide_id,
                    "title": slide.title,
                    "content": slide.content,
                    "image_suggestions": slide.image_suggestions,
                    "layout_type": slide.layout_type,
                    "estimated_duration": slide.estimated_duration
                }
                for slide in slides_package.slides
            ]
        }
        
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
        
        logger.info(f"スライドメタデータ保存完了: {metadata_path}")
